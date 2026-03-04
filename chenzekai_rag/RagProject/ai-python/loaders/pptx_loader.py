"""PPTX loader implementation."""

from __future__ import annotations

from pathlib import Path

from .base import BaseLoader, ParsedDocument


class PptxLoader(BaseLoader):
    """Load PPTX documents into ParsedDocument objects."""

    def load(self, path: str) -> ParsedDocument:
        """Parse a PPTX file and return a ParsedDocument."""
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:  # pragma: no cover - optional dependency
            raise RuntimeError("python-pptx is required for PPTX loading") from exc

        presentation = Presentation(path)
        chunks: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chunks.append(shape.text)
        text = "\n".join(chunks)
        source = Path(path)
        return ParsedDocument(
            doc_id=source.stem,
            text=text.strip(),
            metadata={"source_type": "pptx"},
            source_path=str(source),
        )
